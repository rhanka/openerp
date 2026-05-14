"""Build decision-pack.pptx from structured content.

Run: tools/.venv/bin/python tools/build-decision-pack.py

Output: docs/study/10-mvp-specs/decision-pack.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Palette
NAVY = RGBColor(0x0F, 0x1F, 0x3E)
NAVY_LIGHT = RGBColor(0x1E, 0x3A, 0x8A)
AMBER = RGBColor(0xD9, 0x76, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x05, 0x96, 0x69)
GRAY_900 = RGBColor(0x11, 0x18, 0x27)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_200 = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

STATUS_COLOR = {
    "RESOLVED": GREEN,
    "P0": RED,
    "P1": AMBER,
    "P2": GRAY_500,
}


def make_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=GRAY_900, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_status_pill(slide, x, y, status):
    """Rounded pill with status label."""
    color = STATUS_COLOR.get(status, GRAY_500)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.9), Inches(0.35))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = status
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return pill


def add_header_bar(slide, section):
    add_rect(slide, 0, 0, Inches(13.333), Inches(0.6), NAVY)
    add_text(slide, Inches(0.5), Inches(0.05), Inches(8), Inches(0.5),
             section, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(9.5), Inches(0.05), Inches(3.5), Inches(0.5),
             "MVP Decision Pack — 2026-05-14", size=11, color=GRAY_200,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_page_num(slide, n, total):
    add_text(slide, Inches(12.3), Inches(7.1), Inches(1), Inches(0.3),
             f"{n} / {total}", size=10, color=GRAY_500, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
                 subtitle, size=14, color=GRAY_500)


# ---------- Slides ----------

def slide_title(prs):
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    add_rect(slide, 0, Inches(6.9), prs.slide_width, Inches(0.1), AMBER)
    add_text(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2),
             "MVP Decision Pack", size=60, bold=True, color=WHITE)
    add_text(slide, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.8),
             "Sprint cadrage MVP OpenERP — synthèse arbitrage",
             size=22, color=GRAY_200)
    add_text(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
             "2026-05-14 · Toutes décisions programme RESOLVED", size=14, color=AMBER, bold=True)
    add_text(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
             "6 specs MVP enrichies · double revue Claude + codex · 12 décisions programme arbitrées",
             size=13, color=GRAY_200)
    return slide


def slide_statut(prs):
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Statut")
    add_slide_title(slide, "Statut du sprint")
    y0 = Inches(1.9)
    col_w = Inches(4.15)
    gap = Inches(0.1)

    cols = [
        ("Fait", GREEN, [
            "6 specs MVP enrichies (+1675 lignes)",
            "Double revue Claude + codex (6800 mots)",
            "Rename @entropiq → @sentropic plain MIT (30 docs)",
            "12 décisions programme arbitrées 2026-05-14",
            "Décisions gravées dans 6 specs",
        ]),
        ("À faire", AMBER, [
            "PR @sentropic : MCP + OTel + policy hooks + identité + marketplace + sandbox API",
            "Création shared-entities-v1.md (canon PG-06)",
            "NOTICE racine (Kill Bill, OpenMeter, Superset, Node-RED)",
            "Script CI anti-copy (patterns expression seulement)",
            "Foundation impl (RLS + RBAC + ICU + Money + ApprovalRequest)",
        ]),
        ("Attendu", NAVY_LIGHT, [
            "Spec foundation figée → premier code",
            "@sentropic PR ouverte, autre agent continue dev",
            "Toi en supervision / audit",
            "Pas de nouveau arbitrage requis avant feedback impl",
        ]),
    ]

    for i, (label, color, items) in enumerate(cols):
        x = Inches(0.5) + (col_w + gap) * i
        # Header
        add_rect(slide, x, y0, col_w, Inches(0.55), color)
        add_text(slide, x, y0, col_w, Inches(0.55),
                 label, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_rect(slide, x, y0 + Inches(0.55), col_w, Inches(4.2), GRAY_100)
        for j, item in enumerate(items):
            add_text(slide, x + Inches(0.2), y0 + Inches(0.7) + Inches(0.7) * j,
                     col_w - Inches(0.4), Inches(0.6),
                     "• " + item, size=12, color=GRAY_900)
    return slide


def slide_synthese(prs):
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Synthèse exécutive")
    add_slide_title(slide, "Synthèse exécutive",
                    "Risque licence levé. 5 risques structurels persistants. 12 décisions programme.")

    # Left: 5 risques
    x_left = Inches(0.5)
    y_top = Inches(2.1)
    add_text(slide, x_left, y_top, Inches(6), Inches(0.4),
             "Risques traités par l'arbitrage", size=16, bold=True, color=GREEN)
    risques = [
        "Multi-tenant : UserIdentity + OrganizationMember adopté (PG-02)",
        "Queue : pgmq MVP + JobQueue abstrait pour migration (PG-04)",
        "Activity : préfixage canon (CrmActivity, ProjectTask, ServiceActivity, AgentRun) (PG-06 art.3)",
        "Events : 3 couches séparées AuditEvent / DomainEvent / TimelineEntry (PG-06 art.2)",
        "Anti-copy : owner nommé + script CI tools/anti-copy-grep.sh (PG-12)",
    ]
    for i, r in enumerate(risques):
        add_text(slide, x_left + Inches(0.1), y_top + Inches(0.5) + Inches(0.55) * i,
                 Inches(6.2), Inches(0.5), f"{i+1}.  {r}", size=12, color=GRAY_900)

    # Right: décisions à arbitrer
    x_right = Inches(6.9)
    add_text(slide, x_right, y_top, Inches(6), Inches(0.4),
             "Décisions programme à arbitrer", size=16, bold=True, color=NAVY_LIGHT)
    pgs = [
        ("PG-01", "License @sentropic", "RESOLVED"),
        ("PG-02", "Identité multi-tenant", "RESOLVED"),
        ("PG-03", "Isolation RLS", "RESOLVED"),
        ("PG-04", "Queue engine unifié", "RESOLVED"),
        ("PG-05", "Catalogue i18n", "RESOLVED"),
        ("PG-06", "Canon entités partagées", "RESOLVED"),
        ("PG-07", "ApprovalRequest entité", "RESOLVED"),
        ("PG-08", "Idempotency-Key universel", "RESOLVED"),
        ("PG-09", "Identité agent + délégation", "RESOLVED"),
        ("PG-10", "Stack BI dashboard", "RESOLVED"),
        ("PG-11", "Stack templating (docx + pdf)", "RESOLVED"),
        ("PG-12", "Anti-copy review process", "RESOLVED"),
    ]
    for i, (k, name, status) in enumerate(pgs):
        y = y_top + Inches(0.5) + Inches(0.36) * i
        add_text(slide, x_right + Inches(0.1), y, Inches(0.8), Inches(0.3),
                 k, size=11, bold=True, color=NAVY_LIGHT)
        add_text(slide, x_right + Inches(0.9), y, Inches(4), Inches(0.3),
                 name, size=11, color=GRAY_900)
        add_status_pill(slide, x_right + Inches(5), y - Inches(0.02), status)
    return slide


def slide_section_divider(prs, title, subtitle):
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    add_rect(slide, 0, Inches(3.0), prs.slide_width, Inches(1.5), NAVY_LIGHT)
    add_text(slide, Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.8),
             title, size=44, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.6),
             subtitle, size=18, color=GRAY_200, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, Inches(7.4), prs.slide_width, Inches(0.1), AMBER)
    return slide


def slide_decision(prs, code, name, status, contexte, options, reco, impact):
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Décisions programme")

    # Title row: code + name + status pill
    add_text(slide, Inches(0.5), Inches(0.85), Inches(1.5), Inches(0.6),
             code, size=30, bold=True, color=AMBER)
    add_text(slide, Inches(1.9), Inches(0.85), Inches(9), Inches(0.6),
             name, size=26, bold=True, color=NAVY)
    add_status_pill(slide, Inches(11.7), Inches(1.0), status)

    # Contexte block
    add_text(slide, Inches(0.5), Inches(1.7), Inches(2), Inches(0.3),
             "CONTEXTE", size=11, bold=True, color=GRAY_500)
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(12.3), Inches(0.9), GRAY_100)
    add_text(slide, Inches(0.7), Inches(2.15), Inches(11.9), Inches(0.8),
             contexte, size=13, color=GRAY_900)

    # Options 2-col
    add_text(slide, Inches(0.5), Inches(3.15), Inches(3), Inches(0.3),
             "OPTIONS", size=11, bold=True, color=GRAY_500)
    add_rect(slide, Inches(0.5), Inches(3.5), Inches(7.5), Inches(2.7), GRAY_100)
    for i, opt in enumerate(options):
        y = Inches(3.6) + Inches(0.85) * i
        add_text(slide, Inches(0.7), y, Inches(0.6), Inches(0.4),
                 f"({chr(97+i)})", size=14, bold=True, color=NAVY_LIGHT)
        add_text(slide, Inches(1.2), y, Inches(6.6), Inches(0.7),
                 opt, size=12, color=GRAY_900)

    # Reco
    add_text(slide, Inches(8.3), Inches(3.15), Inches(3), Inches(0.3),
             "RECOMMANDATION", size=11, bold=True, color=AMBER)
    add_rect(slide, Inches(8.3), Inches(3.5), Inches(4.5), Inches(2.7), AMBER)
    add_text(slide, Inches(8.5), Inches(3.65), Inches(4.1), Inches(2.5),
             reco, size=13, bold=True, color=WHITE)

    # Impact bar
    add_rect(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55), NAVY)
    add_text(slide, Inches(0.7), Inches(6.4), Inches(2), Inches(0.55),
             "IMPACT", size=11, bold=True, color=AMBER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(2), Inches(6.4), Inches(10.8), Inches(0.55),
             impact, size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    return slide


def slide_spec_decisions(prs, spec_name, color_label, decisions):
    slide = add_blank_slide(prs)
    add_header_bar(slide, f"Décisions spec — {spec_name}")
    add_slide_title(slide, spec_name, "Décisions clés à graver dans la spec après arbitrage")

    y0 = Inches(2.1)
    add_rect(slide, Inches(0.5), y0, Inches(12.3), Inches(0.5), color_label)
    add_text(slide, Inches(0.7), y0, Inches(1.5), Inches(0.5),
             "REF", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(2.2), y0, Inches(10), Inches(0.5),
             "DÉCISION & RECO", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    for i, (ref, txt) in enumerate(decisions):
        y = y0 + Inches(0.55) + Inches(0.78) * i
        bg = GRAY_100 if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.75), bg)
        add_text(slide, Inches(0.7), y + Inches(0.05), Inches(1.5), Inches(0.7),
                 ref, size=14, bold=True, color=color_label, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(2.2), y + Inches(0.05), Inches(10.5), Inches(0.7),
                 txt, size=12, color=GRAY_900, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_sequence(prs):
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Séquence de décision")
    add_slide_title(slide, "Séquence de décision en 4 phases",
                    "Ordre dans lequel arbitrer pour ne pas bloquer le code")

    phases = [
        ("P0", "Semaine 0 — Bloque code", RED,
         ["PG-01 ✓ résolu", "PG-02 identité multi-tenant", "PG-03 isolation RLS",
          "PG-04 queue", "PG-05 i18n", "PG-06 canon entités"]),
        ("P1", "Semaines 1-2 — Architectural", AMBER,
         ["PG-07 ApprovalRequest", "PG-08 Idempotency-Key", "PG-09 identité agent",
          "PG-10 BI stack", "AuditEvent étendu agentic", "Money type partagé"]),
        ("P2", "Semaines 2-3 — Produit", NAVY_LIGHT,
         ["PG-11 PDF", "PG-12 anti-copy process", "Décisions par spec (5 à 6 par module)",
          "Custom fields declarative", "Timeline grammar", "Supervision UI hybride"]),
        ("P3", "Semaines 3-4 — Opérationnel", GREEN,
         ["NOTICE racine créé", "Script CI anti-copy", "Sweep résiduel terminé",
          "Specs figées prêtes à coder"]),
    ]
    y0 = Inches(2.0)
    h = Inches(1.2)
    for i, (k, title, color, items) in enumerate(phases):
        y = y0 + h * i + Inches(0.05) * i
        # Phase tag
        add_rect(slide, Inches(0.5), y, Inches(1.2), h, color)
        add_text(slide, Inches(0.5), y, Inches(1.2), h,
                 k, size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_rect(slide, Inches(1.8), y, Inches(11), h, GRAY_100)
        add_text(slide, Inches(2.0), y + Inches(0.1), Inches(10.8), Inches(0.4),
                 title, size=14, bold=True, color=color)
        items_str = "  ·  ".join(items)
        add_text(slide, Inches(2.0), y + Inches(0.55), Inches(10.8), Inches(0.6),
                 items_str, size=11, color=GRAY_700)
    return slide


def slide_disagreements(prs):
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Désaccords reviewers")
    add_slide_title(slide, "Points où les reviewers contestent les recos",
                    "À porter à ta connaissance avant arbitrage")

    items = [
        ("Sandbox MVP", "Codex conteste isolated-vm si MVP doit accueillir du code tiers hostile.",
         "Clarifier scope : handlers internes signés seulement → isolated-vm OK."),
        ("Reporting team-scope", "Claude reviewer recommande d'assouplir 'admin-curated only'.",
         "Autoriser dashboards team-scope ; sinon export Excel constant."),
        ("CRM Activity vs Task", "Codex en désaccord avec unification CRM.",
         "Préfère Task ≠ Interaction (cycle de vie distinct)."),
        ("AGT-D-12 fork upstream", "Claude reviewer suggère reformulation neutre.",
         "Éviter engagement irréaliste de contribution upstream best-effort."),
        ("pipeline_stage.{tenant_stage_id}.fr", "Anti-pattern tenant_id dans key i18n.",
         "Table dédiée PipelineStageTranslation au lieu de keys dynamiques."),
        ("Idempotency-Key 4 endpoints", "Reviewers d'accord : trop restrictif.",
         "Étendre à toute POST/DELETE à side-effect métier."),
    ]
    y0 = Inches(2.0)
    for i, (title, issue, reco) in enumerate(items):
        y = y0 + Inches(0.83) * i
        add_rect(slide, Inches(0.5), y, Inches(0.15), Inches(0.7), AMBER)
        add_text(slide, Inches(0.8), y, Inches(3), Inches(0.3),
                 title, size=13, bold=True, color=NAVY)
        add_text(slide, Inches(0.8), y + Inches(0.3), Inches(7), Inches(0.4),
                 issue, size=11, color=GRAY_700)
        add_rect(slide, Inches(8), y, Inches(5), Inches(0.7), GRAY_100)
        add_text(slide, Inches(8.2), y, Inches(4.6), Inches(0.7),
                 "→ " + reco, size=11, color=GRAY_900, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_hotspots(prs):
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Anti-copy hotspots")
    add_slide_title(slide, "Top 12 anti-copy hotspots transverses",
                    "23 hotspots consolidés ; liste complète dans tmp/claude-cross-review.md §9")

    hotspots = [
        ("Odoo mail.thread / Chatter", "LGPL", "CRM, Foundation Comment"),
        ("Twenty ObjectMetadata engine", "AGPL", "CRM custom fields"),
        ("Twenty workspaceMember", "AGPL", "Foundation"),
        ("Odoo ir.model.fields field-ACL", "LGPL", "Foundation"),
        ("Odoo mail.activity / mail.message split", "LGPL", "CRM"),
        ("Odoo l10n_ca", "LGPL", "Billing"),
        ("Odoo account.move / account.move.line", "LGPL", "Billing"),
        ("Odoo auth_totp UI text", "LGPL", "Foundation"),
        ("Odoo ir.rule record rules DSL", "LGPL", "Foundation"),
        ("Kimai timesheet + invoice renderers", "AGPL", "Project"),
        ("Superset chart controls + Selenium PDF", "Apache (défensif)", "Reporting"),
        ("Node-RED node palette + flow JSON", "Apache (défensif)", "Reporting"),
    ]

    y0 = Inches(2.0)
    # Header row
    add_rect(slide, Inches(0.5), y0, Inches(12.3), Inches(0.4), NAVY)
    add_text(slide, Inches(0.7), y0, Inches(0.8), Inches(0.4),
             "#", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.3), y0, Inches(7), Inches(0.4),
             "HOTSPOT", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(8.5), y0, Inches(2), Inches(0.4),
             "DONNEUR", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.5), y0, Inches(2.3), Inches(0.4),
             "SPEC(S) IMPACTÉE(S)", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, (h, lic, spec) in enumerate(hotspots):
        y = y0 + Inches(0.4) + Inches(0.35) * i
        bg = GRAY_100 if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.35), bg)
        add_text(slide, Inches(0.7), y, Inches(0.8), Inches(0.35),
                 str(i+1), size=11, bold=True, color=AMBER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.3), y, Inches(7), Inches(0.35),
                 h, size=11, color=GRAY_900, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(8.5), y, Inches(2), Inches(0.35),
                 lic, size=10, color=GRAY_700, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(10.5), y, Inches(2.3), Inches(0.35),
                 spec, size=10, color=GRAY_700, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_action(prs):
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    add_text(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.6),
             "Action porteur produit", size=36, bold=True, color=AMBER)

    steps = [
        ("1", "Arbitrage terminé", "12 décisions programme arbitrées 2026-05-14. Décisions gravées inline dans les 6 specs MVP."),
        ("2", "PR @sentropic", "Ouvrir PR via ~/src/entropiq pour MCP + OTel + policy hooks + multi-tenant identity + marketplace primitives + sandbox API. Respect du plan @sentropic. L'autre agent continue le dev."),
        ("3", "Implémentation foundation", "Démarrer impl foundation (UserIdentity + OrganizationMember + RLS + RBAC + ICU + Money + ApprovalRequest + Idempotency-Key) + shared-entities-v1.md."),
        ("4", "Mode supervision", "Toi en supervision / audit. Je suis en réalisation. Push à ton OK explicite."),
    ]
    y0 = Inches(1.9)
    for i, (num, title, body) in enumerate(steps):
        y = y0 + Inches(1.15) * i
        add_rect(slide, Inches(0.8), y, Inches(0.9), Inches(0.9), AMBER)
        add_text(slide, Inches(0.8), y, Inches(0.9), Inches(0.9),
                 num, size=40, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.4),
                 title, size=18, bold=True, color=WHITE)
        add_text(slide, Inches(2.0), y + Inches(0.45), Inches(10.5), Inches(0.5),
                 body, size=13, color=GRAY_200)

    add_rect(slide, 0, Inches(7.3), prs.slide_width, Inches(0.2), AMBER)
    return slide


# ---------- Build ----------

def build():
    prs = make_prs()

    # 1 — Title
    slide_title(prs)
    # 2 — Statut
    slide_statut(prs)
    # 3 — Synthèse
    slide_synthese(prs)

    # 4 — Section divider décisions programme
    slide_section_divider(prs, "Décisions programme", "12 décisions transverses · toutes arbitrées 2026-05-14")

    # 5..16 — 12 décisions
    decisions = [
        ("PG-01", "License @sentropic", "RESOLVED",
         "Renommage @entropiq → @sentropic 2026-05-12, license simplifiée en plain MIT à toi-même. AGT-D-12 et AGT-D-13 résolus, plus de blocant commercial.",
         ["Plain MIT (adopté)", "MIT with Commercial Restrictions (rejeté)", "Waiver écrit (non nécessaire)"],
         "Plain MIT adopté. Tout déploiement commercial OpenERP est juridiquement clear côté runtime.",
         "Licence: aucun risque · Timeline: 0j · Action restante: sweep rename terminé (commit 54dede9)"),

        ("PG-02", "Identité + membership multi-tenant", "RESOLVED",
         "User.organization_id (foundation actuel) casse le multi-tenant quand un humain est membre de plusieurs orgs (consultants, partenaires).",
         ["UserIdentity global + OrganizationMember par tenant (codex reco)",
          "Rester User.organization_id mono-org MVP",
          "Multi-org via table UserOrgLink"],
         "✓ Arbitré 2026-05-14 : option (a) UserIdentity + OrganizationMember. Multi-tenant dès MVP, refactor FK et JWT en conséquence.",
         "Foundation P0 · Change toutes FK et tous JWT · Bloque tout schéma métier"),

        ("PG-03", "Isolation multi-tenant", "RESOLVED",
         "Foundation propose row-level organization_id + Postgres RLS. À acter formellement, bloque tous schémas et toutes migrations.",
         ["Row-level + Postgres RLS (foundation reco)",
          "Schema-per-tenant (migrations énormes)",
          "DB-per-tenant (ops lourd MVP)"],
         "✓ Arbitré 2026-05-14 : (a) RLS row-level MVP, avec abstraction TenantIsolationStrategy supportant (a/b/c) dès le contrat foundation. Switch par config sans réécriture downstream.",
         "Foundation P0 · Impacte toutes migrations · Reporting + Agentic dépendent de cette résolution"),

        ("PG-04", "Queue engine unifié", "RESOLVED",
         "Billing propose pgmq. Reporting (scheduled delivery) et Agentic (agent runs) muets. Risque : 3 queue engines coexistants.",
         ["pgmq (Postgres-native, simple ops)",
          "BullMQ + Redis (Node mature)",
          "NATS JetStream (polyglot)",
          "Inngest managed (rejeté MVP)"],
         "✓ Arbitré 2026-05-14 : (a) pgmq MVP, avec interface JobQueue abstraite supportant (a/b/c). Adapter pattern, scheduler ne touche jamais pgmq directement.",
         "Touche billing récurrent + reporting scheduled + agentic runs · 3 implémentations possibles via le même contrat"),

        ("PG-05", "Catalogue i18n unifié", "RESOLVED",
         "Foundation propose ICU JSON nested. Project utilise fr_label/en_label colonnes SQL (anti-pattern). Incohérence à corriger.",
         ["ICU JSON nested + table TranslationKey foundation",
          "gettext .po per-module (anti-copy hotspot Odoo)",
          "Flat JSON simple"],
         "ICU JSON nested + table TranslationKey foundation. Project corrige ServiceActivity pour utiliser keys au lieu de colonnes.",
         "Bloque build pipeline · Foundation owns + 5 specs alignées · FR-CA prioritaire EN-CA parité"),

        ("PG-06", "Canon entités partagées (4 articles)", "RESOLVED",
         "Codex propose un canon transverse cohérent à valider en bloc. Ce ne sont PAS des options alternatives mais les 4 articles du canon à graver. Sans ce canon, chaque module réinvente.",
         ["Article 1 : Organization (métier) + tenant (runtime), Money = {amount_minor, currency(3), scale}",
          "Article 2 : AuditEvent / DomainEvent / TimelineEntry séparés (3 couches)",
          "Article 3 : Activity préfixé : CrmActivity, ProjectTask, ServiceActivity, AgentRun",
          "Article 4 : Company entité + customer = rôle ; Project owns InvoiceProposal, Billing owns Invoice"],
         "✓ Arbitré 2026-05-14 : canon adopté en bloc. Création docs/study/10-mvp-specs/shared-entities-v1.md comme contrat transverse à la prochaine étape.",
         "Contrat transverse · Bloque toutes specs · Empêche divergence terminologique pré-code"),

        ("PG-07", "ApprovalRequest entité foundation", "RESOLVED",
         "Foundation a action approve, project a TimeApproval, CRM a threshold, billing a void/write-off, agentic a approval-in-the-loop. 5 réinventions sans canon.",
         ["Chaque module réinvente (état actuel)",
          "Entité ApprovalRequest foundation partagée",
          "Workflow primitive externalisée vers reporting/automation"],
         "✓ Arbitré 2026-05-14 : (b) ApprovalRequest foundation, exposée via 3 surfaces — API REST, tool MCP request_approval(), SDK @sentropic/openerp-sdk. Codex et Claude consomment le même tool schema.",
         "Touche 5 specs · Économise 4 réimplémentations · Audit trail unifié"),

        ("PG-08", "Idempotency-Key universel", "RESOLVED",
         "Header HTTP Idempotency-Key envoyé par le client sur POST/DELETE. Serveur stocke (key, response) en DB ; replay → cache hit. Évite double facturation, double agent run, double import sur retry réseau.",
         ["4 endpoints seulement (foundation actuel)",
          "Toute POST/DELETE à side-effect métier (extension)",
          "Aucune idempotence MVP (rejeté)"],
         "✓ Arbitré 2026-05-14 : (b) Idempotency-Key obligatoire sur toute POST/DELETE à side-effect métier. TTL 24h, table IdempotencyRecord en foundation. Middleware pattern (Stripe).",
         "Touche toutes specs API · Sécurise retry agentic + automation + import CRM"),

        ("PG-09", "Identité agent + délégation", "RESOLVED",
         "JWT seul insuffisant. Délégation d'action et policy association = objectif clé pour enabler l'agentique. Audit-grade traceabilité requise.",
         ["MVP : JWT + RFC 8693 token exchange (act / may_act claims)",
          "Post-MVP : SPIFFE/SVID workload attestation (CNCF, aligné NIST SP 800-204)",
          "Policy binding : AgentDefinition.allowed_tools, max_spend, requires_approval_above"],
         "✓ Arbitré 2026-05-14 : stack hybride. MVP RFC 8693 token exchange + ApprovalRequest comme issuance. SPIFFE post-MVP via abstraction IdentityProvider. Cookie humains + JWT signé agents. AuditEvent étendu avec acting_principal + on_behalf_of + policy_decision_id.",
         "Foundation + Agentic décident ensemble · Bloque tout flow d'agent autonome · Chaîne délégation auditable de bout en bout"),

        ("PG-10", "Stack BI / dashboard", "RESOLVED",
         "Reporting propose native SvelteKit dashboard. CRM, project, billing ont besoin de dashboards sans engagement clair → risque dérive Superset-clone.",
         ["Native SvelteKit + primitive foundation Widget/Dashboard",
          "Superset embed iframe",
          "Grafana Apache embed (option pack)"],
         "✓ Arbitré 2026-05-14 : (a) primitive Widget foundation, coordination tokens/composants avec @sent-tech à clarifier. Veille charts Svelte (layerchart, svelte-echarts, chart.js, observable plot) à dispatcher. Aspiration VizQL post-MVP.",
         "Touche reporting + CRM dashboards + project dashboards + billing dashboards · Évite dépendance Selenium/Chromium"),

        ("PG-11", "Stack templating (docx + pdf)", "RESOLVED",
         "Reporting propose server-side template ; billing parle de PDF FR/EN sans préciser ; foundation muet. @sentropic a déjà un moteur DOCX mature (dolanmiu/docx + couche maison patchDocument, syntaxe {{...}} + boucles, Word TOC) et un PDF plus basique.",
         ["Extraction @sentropic : librariser docx + pdf, OpenERP consomme",
          "Vivliostyle (CSS Paged Media, Node) — rejeté au profit de @sentropic",
          "LaTeX — rejeté pour stack TS"],
         "✓ Arbitré 2026-05-14 : extraction @sentropic en 2 packages — @sentropic/docx-templating (extractible immédiatement) + @sentropic/pdf-templating (audit + roadmap : Paged Media, FR-CA fonts, accessibilité, templates fiscaux statutaires). OpenERP consomme via npm.",
         "Touche reporting exports + billing invoices + foundation receipts · Coordination avec roadmap @sentropic"),

        ("PG-12", "Anti-copy review process", "RESOLVED",
         "23 hotspots consolidés. Aucun owner. Aucun script CI. Précision : noms business génériques (Company, Invoice, Project, Task, Currency) ne sont pas anti-copy — une pomme est une pomme, non copyrightable.",
         ["Revue spec par spec (atomisé, état actuel)",
          "Owner anti-copy nommé + revue par PR",
          "+ Script CI grep patterns expression-spécifiques (Odoo mail.thread, Twenty ObjectMetadata, etc.)"],
         "✓ Arbitré 2026-05-14 : (b) owner anti-copy + (c) tools/anti-copy-grep.sh ciblé patterns expression seulement (pas les noms génériques business). Liste précise : Odoo ir.*, account.move, mail.thread + Twenty ObjectMetadata + Frappe doctype + Kimai KEvent + SuiteCRM bean->.",
         "Évite contamination LGPL/AGPL · Discipline pre-code obligatoire"),
    ]
    for d in decisions:
        slide_decision(prs, *d)

    # Section divider décisions par spec
    slide_section_divider(prs, "Décisions par spec", "30 décisions clés réparties sur 6 modules MVP")

    # Per-spec slides
    slide_spec_decisions(prs, "Foundation — Security & i18n", NAVY_LIGHT, [
        ("F-1", "RBAC plat resource.action.scope confirmé"),
        ("F-2", "Étendre AuditEvent foundation avec colonnes optionnelles agentiques avant que les autres specs ne dépendent du schéma"),
        ("F-3", "Audit log = table append-only partitionnée par mois, rétention 1 an défaut, plancher 90j"),
        ("F-4", "Auth MVP = passkey / WebAuthn direct (drop password + TOTP MVP)"),
        ("F-5", "Permission granularity = objet uniquement MVP, hook record-policy réservé post-MVP"),
    ])
    slide_spec_decisions(prs, "CRM — Customer Timeline", NAVY_LIGHT, [
        ("C-1", "Lead / Opportunity séparés (Odoo-style en abstraction, pas verbatim)"),
        ("C-2", "Custom fields = JSON column + CustomFieldDefinition declarative limitée (pas de runtime metadata engine)"),
        ("C-3", "Multi-currency par opportunité — à valider avec Billing"),
        ("C-4", "QuoteHandoff : CRM owns handoff event, Billing owns quote doc"),
        ("C-5", "Pipeline stage translation = table dédiée (pas keys dynamiques avec tenant_id)"),
    ])
    slide_spec_decisions(prs, "Project — Time to Invoice", NAVY_LIGHT, [
        ("P-1", "Time entry = manual + timer optionnel (codex réserve : timer post-MVP)"),
        ("P-2", "Approval per-week (per-entry trop granulaire, per-project trop coarse)"),
        ("P-3", "Invoice trigger = mix manuel + milestone. Récurrent post-MVP."),
        ("P-4", "Tâche hiérarchie = parent-enfant 1 niveau MVP. Portfolio post-MVP."),
        ("P-5", "Lien projet ↔ opportunité = 1:N (1 projet peut consommer plusieurs opportunités)"),
    ])
    slide_spec_decisions(prs, "Billing & Accounting", NAVY_LIGHT, [
        ("B-1", "Subscription = native simple + interface compatible Kill Bill (Apache) pour migration"),
        ("B-2", "Plan comptable = pré-chargé CA/QC abstrait + révision juridique avant pilote"),
        ("B-3", "Tax engine = pluggable interne versionné (GST/QST règles)"),
        ("B-4", "Multi-currency = MVP. Book tenant single, transactions multidevises"),
        ("B-5", "Payment = Stripe seul MVP. Multi-PSP via abstraction prête"),
    ])
    slide_spec_decisions(prs, "Reporting & Automation", NAVY_LIGHT, [
        ("R-1", "BI = native SvelteKit + primitive foundation Widget (cf. PG-10)"),
        ("R-2", "Automation = typed YAML DSL + TS handlers. Pas de Node-RED canvas MVP."),
        ("R-3", "Scheduled reports = pgmq cron (cf. PG-04)"),
        ("R-4", "Drill-down = live SQL MVP. Vues matérialisées post-MVP."),
        ("R-5", "Permissions = foundation RBAC + visibility=team MVP. Rôle 'performance opérationnelle' (process défini, KPIs définis, proche mesure marché). Agent assistant aide à cadrer KPIs."),
    ])
    slide_spec_decisions(prs, "Agentic Impacts (transverse)", AMBER, [
        ("A-1", "Sandbox = responsabilité @sentropic. OpenERP exige API sandbox + capability manifest. Feature request si manquant."),
        ("A-2", "Policy = native TS MVP. OPA en pack post-MVP."),
        ("A-3", "Observability = OpenInference + OpenLLMetry MVP"),
        ("A-4", "Identity agent = RFC 8693 token exchange MVP + SPIFFE post-MVP (cf. PG-09)"),
        ("A-5", "MCP registry = interne native MVP. Public registry post-MVP."),
        ("A-6", "Supervision UI = hybride inline + panneau /supervision. Contrainte transverse : tout écran métier réserve zone supervision banner."),
    ])

    # Sequence
    slide_sequence(prs)

    # Disagreements
    slide_disagreements(prs)

    # Hotspots
    slide_hotspots(prs)

    # Action
    slide_action(prs)

    # Add page numbers (skip title + section dividers)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        # Title slide and section dividers don't get page numbers
        if i in (1, 4, 17):
            continue
        add_page_num(slide, i, total)

    prs.save("docs/study/10-mvp-specs/decision-pack.pptx")
    print(f"Generated {total} slides.")


if __name__ == "__main__":
    build()
